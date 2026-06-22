#!/usr/bin/env python3
"""Generate BioVault Sentinel MVP overview deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "BioVault-Sentinel-MVP-Overview.pptx"

DARK = RGBColor(0x0F, 0x17, 0x2A)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
SLATE = RGBColor(0xE2, 0xE8, 0xF0)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title(slide, title: str, subtitle: str = "") -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.2), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(8)


def add_bullets(slide, items: list[str], top: float = 1.6, accent: RGBColor = GREEN) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.5), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = SLATE
        p.space_after = Pt(10)
        p.bullet = True


def add_flow_box(slide, text: str, left: float, top: float, width: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER


def add_arrow(slide, left: float, top: float) -> None:
    slide.shapes.add_shape(25, Inches(left), Inches(top), Inches(0.35), Inches(0.25)).line.color.rgb = GRAY


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Title
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, DARK)
    t = s1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "BioVault Sentinel"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    p2.text = "Phase 1 MVP · v0.1.0-mvp"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "Managed LTO-9 vault + dashboard · Metadata online, tapes air-gapped"
    p3.font.size = Pt(16)
    p3.font.color.rgb = SLATE
    p3.space_before = Pt(20)

    # Slide 2 — Problem / Solution
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, DARK)
    add_title(s2, "Why Sentinel?", "DPDPA 2023 compliance without cloud egress risk")
    add_bullets(
        s2,
        [
            "Hospitals & enterprises must retain records — but on-prem is risky, cloud is expensive",
            "Sentinel: write once to magnetic tape, shelf offline (air-gapped), index searchable online",
            "Clients search & request files; technicians retrieve in 15-minute SLA",
            "Signed PDF certificates for ingest, deletion, and audit export",
        ],
    )

    # Slide 3 — Architecture
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, DARK)
    add_title(s3, "Architecture", "Two portals · one API · metadata DB · tape offline")
    add_flow_box(s3, "Client Portal :5173", 0.6, 1.7, 2.4, GREEN)
    add_flow_box(s3, "Admin Portal :5174", 0.6, 2.6, 2.4, AMBER)
    add_arrow(s3, 3.05, 2.0)
    add_arrow(s3, 3.05, 2.85)
    add_flow_box(s3, "Sentinel API :4000", 3.5, 2.15, 2.8, SLATE)
    add_arrow(s3, 6.4, 2.45)
    add_flow_box(s3, "MongoDB (metadata)", 6.85, 1.7, 2.6, RGBColor(0x60, 0xA5, 0xFA))
    add_flow_box(s3, "Redis + BullMQ", 6.85, 2.6, 2.6, RGBColor(0x60, 0xA5, 0xFA))
    add_arrow(s3, 9.5, 2.45)
    add_flow_box(s3, "LTO-9 Tape (sim)", 9.95, 2.15, 2.6, RGBColor(0xA7, 0x8B, 0xFA))
    note = s3.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(11.5), Inches(2.5))
    ntf = note.text_frame
    np = ntf.paragraphs[0]
    np.text = "Rule: No file bytes in MongoDB or public API — only metadata & transient staging"
    np.font.size = Pt(16)
    np.font.color.rgb = GRAY

    # Slide 4 — Ingest flow
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4, DARK)
    add_title(s4, "Flow 1 — Ingest", "Upload → index → tape write → verify → seal")
    y = 1.75
    steps = [
        ("1. Upload", "Client uploads via portal (multipart)"),
        ("2. Index", "SHA-256 checksum + metadata → MongoDB"),
        ("3. Tape write", "BullMQ worker writes to LTO-9 simulator"),
        ("4. Verify", "Read-back checksum must match"),
        ("5. Seal", "Job sealed · staging purged · signed ingest PDF"),
    ]
    for label, desc in steps:
        lb = s4.shapes.add_textbox(Inches(0.75), Inches(y), Inches(2), Inches(0.4))
        lp = lb.text_frame.paragraphs[0]
        lp.text = label
        lp.font.bold = True
        lp.font.size = Pt(17)
        lp.font.color.rgb = GREEN
        db = s4.shapes.add_textbox(Inches(2.8), Inches(y), Inches(9.5), Inches(0.4))
        dp = db.text_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(17)
        dp.font.color.rgb = SLATE
        y += 0.55

    # Slide 5 — Retrieval flow
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5, DARK)
    add_title(s5, "Flow 2 — Retrieval", "15-minute SLA · admin sees tape location · client does not")
    add_bullets(
        s5,
        [
            "Client searches archive → requests file → job enters queue (dueAt = +15 min)",
            "Alert if unassigned > 60 seconds (audit log)",
            "Technician: Assign → Start → Complete (barcode, rack, slot shown in admin only)",
            "File staged temporarily → client gets expiring download link → purge after use",
        ],
        top=1.65,
        accent=AMBER,
    )

    # Slide 6 — Features
    s6 = prs.slides.add_slide(blank)
    set_slide_bg(s6, DARK)
    add_title(s6, "Live Features (MVP)", "Days 1–16 complete")
    add_bullets(
        s6,
        [
            "Auth & RBAC — separate client / ops sessions (5173 + 5174 simultaneously)",
            "Ingest pipeline — sim tape, verify, ingest confirmation PDF",
            "Search & retrieval — client UI + SLA countdown tracker",
            "Admin ops — job queue, tape inventory (health green/amber/red)",
            "Audit export — immutable hash-chained events",
            "Billing dashboard — storage TB, retrieval bundle, cloud savings estimate",
            "DPDPA erasure — request → degauss → signed deletion certificate",
        ],
        top=1.55,
    )

    # Slide 7 — Demo today
    s7 = prs.slides.add_slide(blank)
    set_slide_bg(s7, DARK)
    add_title(s7, "Demo Data (seeded)", "Run: pnpm --filter @biovault/sentinel-api db:seed")
    add_bullets(
        s7,
        [
            "12 archived files — MRI, CT, lab PDFs (Acme Hospital)",
            "5 retrieval jobs — pending, assigned, in progress, delivered",
            "3 active tapes + inventory health scores",
            "2 erasure requests — 1 awaiting degauss, 1 completed",
            "",
            "Logins: admin@acme.test / tech@biovault.test — ChangeMe123!",
            "Client :5173 · Admin :5174 · API :4000",
        ],
        top=1.55,
    )

    # Slide 8 — Out of scope
    s8 = prs.slides.add_slide(blank)
    set_slide_bg(s8, DARK)
    add_title(s8, "Not in MVP", "Phase 1 completion backlog")
    add_bullets(
        s8,
        [
            "Production SFTP ingest · robotic tape picker",
            "Full MFA · onboarding wizard · production billing/invoicing",
            "Automated tape re-copy · SMS/email alerts",
            "Owning a data centre — colo lease only (CtrlS Hyderabad)",
        ],
        top=1.65,
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
